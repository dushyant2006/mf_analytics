#!/usr/bin/env python3
"""
create_presentation.py
======================
Generates Bluestock_MF_Presentation.pptx (12 slides)
Usage: python create_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

ROOT    = Path(__file__).resolve().parent
REPORTS = ROOT / "reports"
CHARTS  = REPORTS / "charts"
OUT     = REPORTS / "Bluestock_MF_Presentation.pptx"

# Colors
NAVY  = RGBColor(0x0F, 0x1B, 0x2D)
TEAL  = RGBColor(0x00, 0xD4, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0xCC, 0xCC, 0xCC)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def text_box(slide, text, l, t, w, h, size=24, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def accent_bar(slide, t=1.1):
    bar = slide.shapes.add_shape(
        1, Inches(0.5), Inches(t), Inches(12.33), Inches(0.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

def add_image(slide, path, l, t, w, h):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(l), Inches(t),
                                  Inches(w), Inches(h))

# ── Slide 1: Title ──────────────────────────────────────────
s = add_slide(); bg(s)
text_box(s, "🏦 Bluestock Fintech", 0.5, 0.4, 12, 1,
         size=14, color=TEAL)
text_box(s, "Mutual Fund Analytics Platform",
         0.5, 1.2, 12, 1.5, size=40, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
accent_bar(s, 2.9)
text_box(s, "Capstone Project  |  Individual Submission  |  2026",
         0.5, 3.1, 12, 0.6, size=18, color=GRAY, align=PP_ALIGN.CENTER)
text_box(s, "Dushyant Sharma", 0.5, 4.0, 12, 0.6,
         size=16, color=TEAL, align=PP_ALIGN.CENTER)
text_box(s, "Python · SQLite · Power BI · mfapi.in · AMFI Data",
         0.5, 4.7, 12, 0.6, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ── Slide 2: Problem & Objective ────────────────────────────
s = add_slide(); bg(s)
text_box(s, "Problem & Objective", 0.5, 0.3, 12, 0.8,
         size=28, bold=True, color=TEAL)
accent_bar(s)
points = [
    "❓  India's MF industry manages ₹67+ lakh crore AUM — yet retail investors lack easy analytical tools",
    "📊  Objective: Build a full-stack analytics platform from raw AMFI data to interactive dashboard",
    "🎯  Cover ETL → Cleaning → SQL DB → EDA → Risk Metrics → Power BI Dashboard",
    "👥  Target: Retail investors, financial advisors, and fintech analysts",
    "✅  All data from public sources: AMFI India + mfapi.in REST API — no API keys required",
]
for i, pt in enumerate(points):
    text_box(s, pt, 0.7, 1.3 + i*1.0, 12, 0.8, size=15, color=WHITE)

# ── Slide 3: Data Sources ────────────────────────────────────
s = add_slide(); bg(s)
text_box(s, "Data Sources & Scale", 0.5, 0.3, 12, 0.8,
         size=28, bold=True, color=TEAL)
accent_bar(s)
rows = [
    ("📁 Dataset", "Rows", "Description"),
    ("01_fund_master", "40", "40 AMFI schemes — codes, AMC, risk grade"),
    ("02_nav_history", "46,000", "Daily NAV Jan 2022–May 2026"),
    ("03_aum_by_fund_house", "90", "Quarterly AUM for top 10 AMCs"),
    ("04_monthly_sip_inflows", "48", "SIP inflows — ₹31,002 Cr milestone"),
    ("08_investor_transactions", "32,778", "SIP/Lumpsum/Redemption across 12 states"),
    ("10_benchmark_indices", "8,050", "Nifty 50, Nifty 100, BSE SmallCap"),
]
for i, (a, b, c) in enumerate(rows):
    clr = TEAL if i == 0 else WHITE
    sz  = 13 if i == 0 else 12
    text_box(s, a, 0.5, 1.2 + i*0.72, 4.5, 0.6, size=sz, color=clr, bold=(i==0))
    text_box(s, b, 5.1, 1.2 + i*0.72, 1.5, 0.6, size=sz, color=clr, bold=(i==0))
    text_box(s, c, 6.8, 1.2 + i*0.72, 6.0, 0.6, size=sz, color=clr, bold=(i==0))

# ── Slide 4: Architecture ────────────────────────────────────
s = add_slide(); bg(s)
text_box(s, "Data Pipeline & Architecture", 0.5, 0.3, 12, 0.8,
         size=28, bold=True, color=TEAL)
accent_bar(s)
steps = [
    ("1️⃣  EXTRACT",  "mfapi.in REST API + 10 AMFI CSVs"),
    ("2️⃣  TRANSFORM","Pandas cleaning · forward-fill NAV · validate types"),
    ("3️⃣  LOAD",      "SQLite star schema · 11 tables · 5.5 MB DB"),
    ("4️⃣  ANALYSE",  "Sharpe · Sortino · Alpha · Beta · VaR · HHI"),
    ("5️⃣  VISUALISE","Power BI 4-page dashboard · dark fintech theme"),
]
for i, (title, desc) in enumerate(steps):
    text_box(s, title, 0.7, 1.3 + i*1.0, 3.5, 0.7,
             size=15, bold=True, color=TEAL)
    text_box(s, desc,  4.3, 1.3 + i*1.0, 8.5, 0.7,
             size=14, color=WHITE)

# ── Slide 5: EDA Highlights 1 ───────────────────────────────
s = add_slide(); bg(s)
text_box(s, "EDA Highlights — Market Overview", 0.5, 0.3, 12, 0.8,
         size=28, bold=True, color=TEAL)
accent_bar(s)
add_image(s, CHARTS / "03_sip_trend.png",      0.3, 1.2, 6.0, 3.2)
add_image(s, CHARTS / "07_folio_growth.png",   6.6, 1.2, 6.0, 3.2)
text_box(s, "SIP inflows hit ₹31,002 Cr all-time high (Dec 2025)",
         0.3, 4.5, 6.0, 0.6, size=12, color=GRAY)
text_box(s, "Folios doubled: 13.26 Cr → 26.12 Cr in 4 years",
         6.6, 4.5, 6.0, 0.6, size=12, color=GRAY)

# ── Slide 6: EDA Highlights 2 ───────────────────────────────
s = add_slide(); bg(s)
text_box(s, "EDA Highlights — Fund & Investor Analysis", 0.5, 0.3, 12, 0.8,
         size=28, bold=True, color=TEAL)
accent_bar(s)
add_image(s, CHARTS / "09_sector_donut.png",   0.3, 1.2, 6.0, 3.5)
add_image(s, CHARTS / "05_demographics.png",   6.6, 1.2, 6.0, 3.5)
text_box(s, "Banking dominates sector allocation (₹62,840 Cr)",
         0.3, 4.8, 6.0, 0.6, size=12, color=GRAY)
text_box(s, "26–35 age group drives highest SIP volumes",
         6.6, 4.8, 6.0, 0.6, size=12, color=GRAY)

# ── Slide 7: Performance Metrics 1 ──────────────────────────
s = add_slide(); bg(s)
text_box(s, "Performance Metrics — Risk & Return", 0.5, 0.3, 12, 0.8,
         size=28, bold=True, color=TEAL)
accent_bar(s)
add_image(s, CHARTS / "12_sharpe_ratio.png",   0.3, 1.2, 6.0, 3.5)
add_image(s, CHARTS / "13_max_drawdown.png",   6.6, 1.2, 6.0, 3.5)
text_box(s, "Top funds maintain Sharpe > 1.0 consistently",
         0.3, 4.8, 6.0, 0.6, size=12, color=GRAY)
text_box(s, "Small cap funds show worst max drawdown (−24%+)",
         6.6, 4.8, 6.0, 0.6, size=12, color=GRAY)

# ── Slide 8: Performance Metrics 2 ──────────────────────────
s = add_slide(); bg(s)
text_box(s, "Performance Metrics — Scorecard & Benchmark", 0.5, 0.3, 12, 0.8,
         size=28, bold=True, color=TEAL)
accent_bar(s)
add_image(s, CHARTS / "14_fund_scorecard.png",      0.3, 1.2, 6.0, 3.5)
add_image(s, CHARTS / "15_benchmark_comparison.png",6.6, 1.2, 6.0, 3.5)
text_box(s, "Composite scorecard ranks all 40 funds (0–100)",
         0.3, 4.8, 6.0, 0.6, size=12, color=GRAY)
text_box(s, "Top funds outperform Nifty 50 on normalised basis",
         6.6, 4.8, 6.0, 0.6, size=12, color=GRAY)

# ── Slide 9: Dashboard Screenshot 1 ─────────────────────────
s = add_slide(); bg(s)
text_box(s, "Power BI Dashboard — Industry Overview & Fund Performance",
         0.5, 0.3, 12, 0.8, size=24, bold=True, color=TEAL)
accent_bar(s)
add_image(s, ROOT/"dashboard"/"page1_industry_overview.png",  0.3, 1.2, 6.0, 4.5)
add_image(s, ROOT/"dashboard"/"page2_fund_performance.png",   6.6, 1.2, 6.0, 4.5)

# ── Slide 10: Dashboard Screenshot 2 ────────────────────────
s = add_slide(); bg(s)
text_box(s, "Power BI Dashboard — Investor Analytics & SIP Trends",
         0.5, 0.3, 12, 0.8, size=24, bold=True, color=TEAL)
accent_bar(s)
add_image(s, ROOT/"dashboard"/"page3_investor_analytics.png", 0.3, 1.2, 6.0, 4.5)
add_image(s, ROOT/"dashboard"/"page4_sip_trends.png",         6.6, 1.2, 6.0, 4.5)

# ── Slide 11: Key Findings ───────────────────────────────────
s = add_slide(); bg(s)
text_box(s, "Key Findings & Recommendations", 0.5, 0.3, 12, 0.8,
         size=28, bold=True, color=TEAL)
accent_bar(s)
findings = [
    "🏆  SBI MF leads industry AUM at ₹12.5L Cr — 106% growth since 2022",
    "📈  SIP inflows hit ₹31,002 Cr all-time high in Dec 2025 — retail boom",
    "📊  Folios doubled to 26.12 Cr — massive financial inclusion success",
    "⚠️   Small cap VaR exceeds −2.5%/day — suitable only for 5Y+ horizon",
    "🎯  Top 3 recommended funds (Moderate risk): HDFC Top 100, Mirae Large Cap, ICICI Bluechip",
    "🔍  75%+ investors maintain SIP continuity — strong discipline observed",
]
for i, f in enumerate(findings):
    text_box(s, f, 0.7, 1.3 + i*0.9, 12, 0.75, size=14, color=WHITE)

# ── Slide 12: Thank You ──────────────────────────────────────
s = add_slide(); bg(s)
text_box(s, "Thank You", 0.5, 1.5, 12, 1.5,
         size=48, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
accent_bar(s, 3.2)
text_box(s, "Dushyant Sharma", 0.5, 3.4, 12, 0.7,
         size=20, color=WHITE, align=PP_ALIGN.CENTER)
text_box(s, "Bluestock Fintech Internship Capstone  |  2026",
         0.5, 4.1, 12, 0.6, size=15, color=GRAY, align=PP_ALIGN.CENTER)
text_box(s, "github.com/dushyant2006/mf_analytics",
         0.5, 4.8, 12, 0.6, size=14, color=TEAL, align=PP_ALIGN.CENTER)

prs.save(str(OUT))
print(f"✓ Presentation saved: {OUT}")
print(f"  Slides: {len(prs.slides)}")